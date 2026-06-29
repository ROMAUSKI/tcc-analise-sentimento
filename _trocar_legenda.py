"""Troca SOMENTE a figura da legenda (canto inf. direito) do slide 13.
Nao move nem altera as outras 5 figuras. Backup antes."""
import glob, os, shutil, datetime
from pptx import Presentation
from pptx.util import Inches

BASE = r"C:\Users\Davi\Documentos\tcc-analise-sentimento"
PPTX = glob.glob(os.path.join(BASE, "Apresenta*", "Apresentacao_TCC_Davi.pptx"))[0]
LEG = os.path.join(BASE, "resultados", "slide13_visoes", "legenda.png")

ts = datetime.datetime.now().strftime("%Y%m%d_%H%M%S")
bdir = os.path.join(os.path.dirname(PPTX), "_backups")
os.makedirs(bdir, exist_ok=True)
shutil.copy2(PPTX, os.path.join(bdir, f"Apresentacao_TCC_Davi_backup_{ts}.pptx"))

prs = Presentation(PPTX)
slide = prs.slides[12]

# a legenda e a imagem mais a direita e mais embaixo (left>7, top>4)
legenda = None
for sh in list(slide.shapes):
    if sh.shape_type == 13:
        L = sh.left / 914400; T = sh.top / 914400
        if L > 7 and T > 4:
            legenda = sh
if legenda is None:
    raise SystemExit("Legenda nao localizada no slide 13.")

left, top, width = legenda.left, legenda.top, legenda.width
print(f"[legenda] encontrada em ({round(left/914400,2)}, {round(top/914400,2)}) w={round(width/914400,2)}")
legenda._element.getparent().remove(legenda._element)
pic = slide.shapes.add_picture(LEG, left, top, width=width)
print(f"[trocada] nova legenda inserida na mesma posicao")

prs.save(PPTX)
print(f"[salvo] {PPTX}")
