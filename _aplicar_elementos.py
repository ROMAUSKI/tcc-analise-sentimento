# -*- coding: utf-8 -*-
"""
Adiciona ELEMENTOS VISUAIS ricos (mesmo estilo dos slides 8-13: chevrons,
cards) aos slides so-texto, SEM regenerar o deck e SEM alterar o texto.
Idempotente: remove o que ele mesmo adicionou (name ELEM_AUTO*) antes de
readicionar. Faz backup. Nao toca em outras imagens/shapes.
"""
import glob, os, shutil, datetime
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

BASE = r"C:\Users\Davi\Documentos\tcc-analise-sentimento"
PPTX = glob.glob(os.path.join(BASE, "Apresenta*", "Apresentacao_TCC_Davi.pptx"))[0]
VINHO = RGBColor(0x8C, 0x23, 0x32)
BRANCO = RGBColor(0xFF, 0xFF, 0xFF)
PRETO = RGBColor(0x22, 0x22, 0x22)
CINZA = RGBColor(0xEC, 0xEC, 0xEC)
CINZAESC = RGBColor(0x5A, 0x5A, 0x5A)
TAG = "ELEM_AUTO"


def _tag(sh, n):
    sh.name = f"{TAG}_{n}"


def chevron(slide, left, top, w, h, title, sub, fill, n):
    ch = slide.shapes.add_shape(MSO_SHAPE.CHEVRON, Inches(left), Inches(top), Inches(w), Inches(h))
    ch.fill.solid(); ch.fill.fore_color.rgb = fill
    ch.line.fill.background(); ch.shadow.inherit = False
    tf = ch.text_frame; tf.word_wrap = True; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = title; r.font.bold = True; r.font.size = Pt(13); r.font.color.rgb = BRANCO
    p2 = tf.add_paragraph(); p2.alignment = PP_ALIGN.CENTER
    r2 = p2.add_run(); r2.text = sub; r2.font.size = Pt(9.5); r2.font.color.rgb = BRANCO
    _tag(ch, n)


def _body(slide):
    # o corpo e o placeholder de maior altura (titulo e n. do slide sao menores)
    cands = [sh for sh in slide.shapes if sh.is_placeholder and sh.height]
    return max(cands, key=lambda s: s.height) if cands else None


def mover_texto(slide, top, height):
    sh = _body(slide)
    if sh is not None:
        sh.top = Inches(top)
        sh.height = Inches(height)


def stat_card(slide, left, top, w, h, number, label, n):
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(w), Inches(h))
    box.fill.solid(); box.fill.fore_color.rgb = BRANCO
    box.line.color.rgb = VINHO; box.line.width = Pt(1.5); box.shadow.inherit = False
    tf = box.text_frame; tf.word_wrap = True; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p1 = tf.paragraphs[0]; p1.alignment = PP_ALIGN.CENTER
    r1 = p1.add_run(); r1.text = number; r1.font.size = Pt(40); r1.font.bold = True; r1.font.color.rgb = VINHO
    p2 = tf.add_paragraph(); p2.alignment = PP_ALIGN.CENTER
    r2 = p2.add_run(); r2.text = label; r2.font.size = Pt(12); r2.font.color.rgb = PRETO
    _tag(box, n)


def info_card(slide, left, top, w, h, titulo, desc, n, fill=None):
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(w), Inches(h))
    box.fill.solid()
    if fill is not None:
        box.fill.fore_color.rgb = fill; box.line.fill.background()
        tcol_t, tcol_d = BRANCO, BRANCO
    else:
        box.fill.fore_color.rgb = BRANCO; box.line.color.rgb = VINHO; box.line.width = Pt(1.5)
        tcol_t, tcol_d = VINHO, PRETO
    box.shadow.inherit = False
    tf = box.text_frame; tf.word_wrap = True; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p1 = tf.paragraphs[0]; p1.alignment = PP_ALIGN.CENTER
    r1 = p1.add_run(); r1.text = titulo; r1.font.size = Pt(15); r1.font.bold = True; r1.font.color.rgb = tcol_t
    p2 = tf.add_paragraph(); p2.alignment = PP_ALIGN.CENTER
    r2 = p2.add_run(); r2.text = desc; r2.font.size = Pt(10.5); r2.font.color.rgb = tcol_d
    _tag(box, n)


def chevrons_row(slide, top, items, tag, h=0.9):
    cw, cstep = 3.25, 2.78
    for i, (tt, sub) in enumerate(items):
        fill = VINHO if i % 2 == 0 else CINZAESC
        chevron(slide, 0.55 + i * cstep, top, cw, h, tt, sub, fill, f"{tag}{i}")


# ---------------- construtores por slide ----------------

def conceito_intro(slide):  # slide 3 (NAO usado — Davi ajustou)
    mover_texto(slide, 2.85, 3.85)
    chevrons_row(slide, 1.62, [
        ("LLMs geram", "frases rotuladas"), ("Treina", "classificador"),
        ("Testa", "reviews reais"), ("Funciona?", "cross-domain")], "intro", h=0.95)


def sumario(slide):  # slide 2 — roteiro macro embaixo
    mover_texto(slide, 1.65, 3.8)
    chevrons_row(slide, 5.6, [
        ("Contexto", "intro · objetivos"), ("Método", "5 visões"),
        ("Resultados", "sintético × real"), ("Conclusão", "viabilidade")], "sum")


def objetivos(slide):  # slide 4
    mover_texto(slide, 1.65, 3.75)
    chevrons_row(slide, 5.6, [
        ("Assertividade", "treino sintético"), ("Coerência", "entre LLMs"),
        ("Cross-domain", "sintético → real"), ("Dois nichos", "subjetividade")], "obj")


def justificativa(slide):  # slide 5 — contraste
    mover_texto(slide, 1.65, 2.8)
    info_card(slide, 0.7, 4.6, 5.6, 1.9,
              "Anotação humana", "cara, lenta e escassa em português", "ju0", fill=CINZAESC)
    info_card(slide, 7.05, 4.6, 5.6, 1.9,
              "LLMs", "geram texto já rotulado, sem anotação manual", "ju1", fill=VINHO)


def referencial(slide):  # slide 6 — evolução dos modelos
    mover_texto(slide, 1.65, 3.75)
    chevrons_row(slide, 5.6, [
        ("TF-IDF", "clássicos"), ("LSTM", "rede neural"),
        ("BERTimbau", "pré-treino pt-BR"), ("LLMs", "geração de dados")], "ref")


def materiais(slide):  # slide 7 — 3 categorias
    mover_texto(slide, 1.65, 3.2)
    info_card(slide, 0.46, 5.05, 3.9, 1.45,
              "3 LLMs", "ChatGPT · Gemini · Claude", "mat0")
    info_card(slide, 4.71, 5.05, 3.9, 1.45,
              "5 classificadores", "NB · LR · SVM · LSTM · BERT", "mat1")
    info_card(slide, 8.96, 5.05, 3.9, 1.45,
              "2 nichos reais", "UTLC-Movies · UTLC-Apps", "mat2")


def conclusao_objetivo(slide):  # slide 14 — stat cards
    mover_texto(slide, 1.65, 3.1)
    stat_card(slide, 1.16, 4.95, 5.0, 1.65, "97%", "Sintético (V1, BERTimbau)", "co0")
    stat_card(slide, 7.16, 4.95, 5.0, 1.65, "43–56%", "Cross-domain real (V5)", "co1")


def conclusao_dificuldades(slide):  # slide 15 — 3 cards
    mover_texto(slide, 1.65, 3.2)
    info_card(slide, 0.46, 5.05, 3.9, 1.45,
              "Geração manual", "~3–4 h por nicho, via interface web", "cd0")
    info_card(slide, 4.71, 5.05, 3.9, 1.45,
              "RAM do Colab", "subamostra estratificada de 100 mil", "cd1")
    info_card(slide, 8.96, 5.05, 3.9, 1.45,
              "Paradoxo", "calibrar o sintético exige dado real", "cd2")


def conclusao_futuros(slide):  # slide 16 — 4 chevrons
    mover_texto(slide, 1.65, 3.75)
    chevrons_row(slide, 5.6, [
        ("Adaptação", "domínio: real + sint."), ("Outros nichos", "+ subjetividade"),
        ("Prompts", "refino iterativo"), ("API", "temperatura · + LLMs")], "cf")


def limpar(slide):
    for sh in list(slide.shapes):
        if sh.name.startswith(TAG):
            sh._element.getparent().remove(sh._element)


# idx 0-based -> funcao construtora.  (slide 3 / idx 2 ficou de fora: Davi ajustou)
MAPA = {
    1: sumario,
    3: objetivos,
    4: justificativa,
    5: referencial,
    6: materiais,
    13: conclusao_objetivo,
    14: conclusao_dificuldades,
    15: conclusao_futuros,
}


def main():
    ts = datetime.datetime.now().strftime("%Y%m%d_%H%M%S")
    bdir = os.path.join(os.path.dirname(PPTX), "_backups")
    os.makedirs(bdir, exist_ok=True)
    shutil.copy2(PPTX, os.path.join(bdir, f"Apresentacao_TCC_Davi_backup_{ts}.pptx"))
    prs = Presentation(PPTX)
    for idx, fn in MAPA.items():
        slide = prs.slides[idx]
        limpar(slide)
        fn(slide)
        print(f"[elem] slide {idx+1}: {fn.__name__}")
    prs.save(PPTX)
    print(f"[salvo] {PPTX}")


if __name__ == "__main__":
    main()
