"""
Aplica um icone tematico (canto inferior direito) nos slides definidos em MAPA.
Idempotente: remove icones que ele mesmo adicionou antes (name ICONE_AUTO*)
sem tocar em outras imagens (ex.: as 5 figuras do slide 13).
Faz backup antes de salvar. NAO altera texto de nenhum slide.
"""
import glob, os, shutil, datetime
from pptx import Presentation
from pptx.util import Inches

BASE = r"C:\Users\Davi\Documentos\tcc-analise-sentimento"
PPTX = glob.glob(os.path.join(BASE, "Apresenta*", "Apresentacao_TCC_Davi.pptx"))[0]
ICONES = os.path.join(BASE, "resultados", "icones_slides")

# idx 0-based -> nome do icone (sem .png).  DEMO: so 2 slides.
MAPA = {
    3:  "target",   # slide 4  Objetivos
    13: "check",    # slide 14 Conclusao 1/3
}

ICON_W = 1.1
LEFT = 13.333 - ICON_W - 0.55   # canto inferior direito
TOP = 5.5
TAG = "ICONE_AUTO"


def main():
    ts = datetime.datetime.now().strftime("%Y%m%d_%H%M%S")
    bdir = os.path.join(os.path.dirname(PPTX), "_backups")
    os.makedirs(bdir, exist_ok=True)
    bkp = os.path.join(bdir, f"Apresentacao_TCC_Davi_backup_{ts}.pptx")
    shutil.copy2(PPTX, bkp)
    print(f"[backup] {bkp}")

    prs = Presentation(PPTX)
    for idx, icone in MAPA.items():
        slide = prs.slides[idx]
        # remove icones automaticos antigos (mantem outras imagens)
        for sh in list(slide.shapes):
            if sh.shape_type == 13 and sh.name.startswith(TAG):
                sh._element.getparent().remove(sh._element)
        path = os.path.join(ICONES, f"{icone}.png")
        pic = slide.shapes.add_picture(path, Inches(LEFT), Inches(TOP), width=Inches(ICON_W))
        pic.name = f"{TAG}_{icone}"
        print(f"[icone] slide {idx+1}: {icone}")

    prs.save(PPTX)
    print(f"[salvo] {PPTX}")


if __name__ == "__main__":
    main()
