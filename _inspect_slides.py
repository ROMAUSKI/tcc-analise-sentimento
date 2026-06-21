"""Inspeciona slides-alvo (somente leitura) para planejar onde encaixar icones."""
import glob, os
from pptx import Presentation

BASE = r"C:\Users\Davi\Documentos\tcc-analise-sentimento"
PPTX = glob.glob(os.path.join(BASE, "Apresenta*", "Apresentacao_TCC_Davi.pptx"))[0]
prs = Presentation(PPTX)

alvo = [1, 2, 3, 4, 5, 6, 13, 14, 15, 16, 17]  # slides 2-7 e 14-18 (0-based)
for idx in alvo:
    s = prs.slides[idx]
    print(f"\n=== Slide {idx+1} (idx {idx}) ===")
    for sh in s.shapes:
        def p(v):
            return None if v is None else round(v/914400, 2)
        txt = sh.text_frame.text.replace("\n", " | ")[:55] if sh.has_text_frame else ""
        pic = " [PIC]" if sh.shape_type == 13 else ""
        print(f"  {sh.name!r:22} pos=({p(sh.left)},{p(sh.top)}) "
              f"size=({p(sh.width)},{p(sh.height)}){pic} {txt!r}")
