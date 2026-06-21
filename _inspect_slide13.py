"""Inspeciona o slide 13 do deck para planejar a edição (somente leitura)."""
import glob, os
from pptx import Presentation
from pptx.util import Emu

BASE = r"C:\Users\Davi\Documentos\tcc-analise-sentimento"
PPTX = glob.glob(os.path.join(BASE, "Apresenta*", "Apresentacao_TCC_Davi.pptx"))[0]
print("Arquivo:", PPTX)

prs = Presentation(PPTX)
print("Total de slides:", len(prs.slides))
print("Tamanho slide (EMU):", prs.slide_width, prs.slide_height,
      "=>", round(prs.slide_width/914400, 2), "x", round(prs.slide_height/914400, 2), "pol")

# slide 13 visivel = indice 12
for idx in (11, 12, 13):
    if idx >= len(prs.slides):
        continue
    s = prs.slides[idx]
    print(f"\n=== Slide indice {idx} (visivel {idx+1}) ===")
    for sh in s.shapes:
        kind = sh.shape_type
        txt = ""
        if sh.has_text_frame:
            txt = sh.text_frame.text.replace("\n", " | ")[:70]
        def p(v):
            return None if v is None else round(v/914400, 2)
        print(f"  - id={sh.shape_id} name={sh.name!r} type={kind} "
              f"pos=({p(sh.left)},{p(sh.top)}) size=({p(sh.width)},{p(sh.height)}) "
              f"is_picture={sh.shape_type==13} txt={txt!r}")
