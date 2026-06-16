# -*- coding: utf-8 -*-
"""Gera a apresentacao do TCC (18 slides) a partir do template UTFPR:
texto enxuto + visual (pipeline, cards, reality gap), sem data no rodape,
Metodos em 2 slides, Conclusao em 3, e notas explicativas (de fixacao) em cada slide.
Tambem gera Roteiro_Apresentacao.md com as mesmas notas."""
import os
import glob
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

BASE = os.path.dirname(os.path.abspath(__file__))
# localizar por glob para evitar problemas de normalizacao Unicode (acentos) no Windows
PASTA_APR = [p for p in glob.glob(os.path.join(BASE, "Apresenta*")) if os.path.isdir(p)][0]
SRC = glob.glob(os.path.join(PASTA_APR, "Modelo*Banca*.pptx"))[0]
OUT_PASTA = os.path.join(PASTA_APR, "Apresentacao_TCC_Davi.pptx")
OUT_ROOT = os.path.join(BASE, "Apresentacao_TCC_Davi.pptx")
ROTEIRO = os.path.join(PASTA_APR, "Roteiro_Apresentacao.md")
IMG = os.path.join(BASE, "artigo", "imagens")
VINHO = RGBColor(0x8C, 0x23, 0x32)
BRANCO = RGBColor(0xFF, 0xFF, 0xFF)
PRETO = RGBColor(0x22, 0x22, 0x22)
CINZA = RGBColor(0xEC, 0xEC, 0xEC)
CINZAESC = RGBColor(0x5A, 0x5A, 0x5A)

prs = Presentation(SRC)


def ph(slide, idx):
    for sh in slide.shapes:
        if sh.is_placeholder and sh.placeholder_format.idx == idx:
            return sh
    return None


def set_title(slide, text, size=None):
    tf = ph(slide, 0).text_frame
    p = tf.paragraphs[0]
    for extra in list(tf.paragraphs)[1:]:
        extra._p.getparent().remove(extra._p)
    if p.runs:
        p.runs[0].text = text
        for r in list(p.runs)[1:]:
            r._r.getparent().remove(r._r)
    else:
        p.add_run().text = text
    if size:
        p.runs[0].font.size = Pt(size)


def set_lines(shape, lines, size=None, align=None):
    tf = shape.text_frame
    tf.clear()
    for i, ln in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        r = p.add_run()
        r.text = ln
        if size:
            r.font.size = Pt(size)
        if align is not None:
            p.alignment = align


def set_body(slide, items, size=None):
    tf = ph(slide, 1).text_frame
    tf.clear()
    for i, (txt, lvl) in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.level = lvl
        r = p.add_run()
        r.text = txt
        if size:
            r.font.size = Pt(size)


def shrink_body(slide, left, top, width, height):
    sh = ph(slide, 1)
    sh.left, sh.top, sh.width, sh.height = Inches(left), Inches(top), Inches(width), Inches(height)


def add_img(slide, path, left, top, width):
    return slide.shapes.add_picture(path, Inches(left), Inches(top), width=Inches(width))


def remove_date(slide):
    sh = ph(slide, 10)
    if sh is not None:
        sh._element.getparent().remove(sh._element)


def decorate_title(slide):
    t = ph(slide, 0)
    t.left, t.top, t.width, t.height = Inches(1.18), Inches(0.4), Inches(11.0), Inches(1.45)
    mk = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.55), Inches(0.92), Inches(0.42), Inches(0.42))
    mk.fill.solid()
    mk.fill.fore_color.rgb = VINHO
    mk.line.fill.background()
    mk.shadow.inherit = False


def stat_card(slide, left, top, w, h, number, label):
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(w), Inches(h))
    box.fill.solid()
    box.fill.fore_color.rgb = BRANCO
    box.line.color.rgb = VINHO
    box.line.width = Pt(1.5)
    box.shadow.inherit = False
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p1 = tf.paragraphs[0]
    p1.alignment = PP_ALIGN.CENTER
    r1 = p1.add_run()
    r1.text = number
    r1.font.size = Pt(50)
    r1.font.bold = True
    r1.font.color.rgb = VINHO
    p2 = tf.add_paragraph()
    p2.alignment = PP_ALIGN.CENTER
    r2 = p2.add_run()
    r2.text = label
    r2.font.size = Pt(13)
    r2.font.color.rgb = PRETO


def chevron(slide, left, top, w, h, title, sub, fill):
    ch = slide.shapes.add_shape(MSO_SHAPE.CHEVRON, Inches(left), Inches(top), Inches(w), Inches(h))
    ch.fill.solid()
    ch.fill.fore_color.rgb = fill
    ch.line.fill.background()
    ch.shadow.inherit = False
    tf = ch.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = title
    r.font.bold = True
    r.font.size = Pt(14)
    r.font.color.rgb = BRANCO
    p2 = tf.add_paragraph()
    p2.alignment = PP_ALIGN.CENTER
    r2 = p2.add_run()
    r2.text = sub
    r2.font.size = Pt(10)
    r2.font.color.rgb = BRANCO


def pipeline(slide, top):
    etapas = [
        ("1. Geração", "3 LLMs · 1.800 frases", VINHO),
        ("2. Pré-proc.", "limpeza · TF-IDF", CINZAESC),
        ("3. Treino", "5 classificadores", VINHO),
        ("4. Avaliação", "5 visões · F1-macro", CINZAESC),
    ]
    cw, cstep, h = 3.25, 2.78, 1.0
    for i, (tt, sub, fill) in enumerate(etapas):
        chevron(slide, 0.55 + i * cstep, top, cw, h, tt, sub, fill)


def visoes_table(slide, top):
    t = slide.shapes.add_table(6, 4, Inches(1.85), Inches(top), Inches(9.6), Inches(2.6)).table
    t.columns[0].width = Inches(1.4)
    for c in (1, 2, 3):
        t.columns[c].width = Inches(2.73)
    dados = [
        ["Visão", "Treino", "Teste", "Volume"],
        ["V1", "Sintético", "Sintético", "Controlado"],
        ["V2", "Real", "Real", "Controlado"],
        ["V3", "Sintético", "Real", "Controlado"],
        ["V4", "Real", "Real", "Desbalanceado"],
        ["V5", "Sintético", "Real", "Desbalanceado"],
    ]
    for ri, row in enumerate(dados):
        for ci, val in enumerate(row):
            cell = t.cell(ri, ci)
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            cell.margin_top = cell.margin_bottom = Pt(1)
            cell.text_frame.clear()
            p = cell.text_frame.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER
            r = p.add_run()
            r.text = val
            r.font.size = Pt(14)
            if ri == 0:
                r.font.bold = True
                r.font.color.rgb = BRANCO
                cell.fill.solid()
                cell.fill.fore_color.rgb = VINHO
            else:
                r.font.color.rgb = PRETO
                cell.fill.solid()
                cell.fill.fore_color.rgb = BRANCO if ri % 2 else CINZA


def set_notes(slide, texto):
    slide.notes_slide.notes_text_frame.text = texto


# ===================== ESTRUTURA =====================
(s_capa, s_sum, s_intro, s_obj, s_just, s_refteo, s_mat, s_met1,
 s_r1, s_r2, s_volume, s_conc1, s_refs, s_fim) = list(prs.slides)

layout_conteudo = s_intro.slide_layout
s_met2 = prs.slides.add_slide(layout_conteudo)
s_rgap = prs.slides.add_slide(layout_conteudo)
s_conc2 = prs.slides.add_slide(layout_conteudo)
s_conc3 = prs.slides.add_slide(layout_conteudo)

# reordenar sldIdLst para a sequencia final
sldIdLst = prs.slides._sldIdLst
ids = list(sldIdLst)
# indices atuais: 0capa 1sum 2intro 3obj 4just 5refteo 6mat 7met1 8r1 9r2 10volume
#                 11conc1 12refs 13fim 14met2 15rgap 16conc2 17conc3
ordem = [0, 1, 2, 3, 4, 5, 6, 7, 14, 8, 9, 15, 10, 11, 16, 17, 12, 13]
for e in ids:
    sldIdLst.remove(e)
for i in ordem:
    sldIdLst.append(ids[i])

# remover data de todos
for s in (s_capa, s_sum, s_intro, s_obj, s_just, s_refteo, s_mat, s_met1, s_met2,
          s_r1, s_r2, s_rgap, s_volume, s_conc1, s_conc2, s_conc3, s_refs, s_fim):
    remove_date(s)

# ===================== CONTEUDO =====================
# 1 CAPA
set_title(s_capa, "Utilização de LLMs para Geração de Bases de Treinamento em Classificação de Sentimento: Uma Análise de Viabilidade Prática", size=28)
set_lines(ph(s_capa, 1), ["Trabalho de Conclusão de Curso 2", "Discente: Davi Romauski Meurer", "Orientador: Prof. Marlon Marcon"])

# 2 SUMARIO
decorate_title(s_sum)

# 3 INTRODUCAO
decorate_title(s_intro)
set_body(s_intro, [
    ("Análise de sentimento: classificar opiniões em positiva, negativa ou neutra.", 0),
    ("Treinar modelos exige datasets rotulados — caros e demorados de construir.", 0),
    ("Em português brasileiro, faltam recursos anotados.", 0),
    ("LLMs podem gerar esses dados já rotulados — mas funcionam no mundo real?", 0),
])

# 4 OBJETIVOS
decorate_title(s_obj)
set_body(s_obj, [
    ("Objetivo geral:", 0),
    ("Avaliar a viabilidade prática de treinar classificadores de sentimento com dados sintéticos de LLMs.", 1),
    ("Objetivos específicos:", 0),
    ("Medir a assertividade com treino apenas sintético.", 1),
    ("Verificar a coerência entre os LLMs geradores.", 1),
    ("Avaliar o cross-domain: treino sintético, teste real.", 1),
    ("Comparar dois nichos com subjetividade diferente.", 1),
])

# 5 JUSTIFICATIVA
decorate_title(s_just)
set_body(s_just, [
    ("Datasets rotulados são caros, e o português tem poucos recursos.", 0),
    ("LLMs geram texto rotulado sem anotação humana.", 0),
    ("Três LLMs reduzem o viés de uma fonte única.", 0),
    ("Dois nichos e o cross-domain revelam onde a abordagem funciona.", 0),
])

# 6 REFERENCIAL
decorate_title(s_refteo)
set_body(s_refteo, [
    ("Análise de sentimento = classificação supervisionada de texto (Pang, 2002).", 0),
    ("TF-IDF com classificadores clássicos: Naive Bayes, Regressão Logística e SVM.", 0),
    ("Modelos neurais: LSTM e BERTimbau (BERT pré-treinado em português).", 0),
    ("Geração de dados sintéticos por LLMs: área de pesquisa recente.", 0),
    ("Métrica principal: F1-macro, justo com classes desbalanceadas.", 0),
])

# 7 MATERIAIS
decorate_title(s_mat)
set_body(s_mat, [
    ("LLMs geradores: ChatGPT, Gemini e Claude.", 0),
    ("Classificadores: Naive Bayes, Reg. Logística, SVM, LSTM e BERTimbau.", 0),
    ("Stack: Python, scikit-learn, PyTorch e Hugging Face Transformers.", 0),
    ("Dados reais: UTLC-Movies e UTLC-Apps (Kaggle), via kagglehub.", 0),
    ("Execução dos modelos neurais em Google Colab (GPU NVIDIA T4).", 0),
])

# 8 METODOS 1/2
decorate_title(s_met1)
set_title(s_met1, "Métodos (1/2): Geração e Dados")
pipeline(s_met1, 1.65)
set_body(s_met1, [
    ("Geração: 3 LLMs produziram 1.800 frases sintéticas por nicho (3 × 3 classes × 200).", 0),
    ("Dois nichos: filmes e séries e aplicativos móveis, com prompts isomorfos.", 0),
    ("Dados reais (cross-domain): UTLC-Movies e UTLC-Apps, ~100 mil reviews/nicho; nota 1–5 mapeada para as 3 classes.", 0),
])
shrink_body(s_met1, 1.18, 3.05, 11.2, 3.0)

# 9 METODOS 2/2
decorate_title(s_met2)
set_title(s_met2, "Métodos (2/2): Avaliação")
set_body(s_met2, [
    ("Pré-processamento: minúsculas, remoção de pontuação e deduplicação.", 0),
    ("Vetorização TF-IDF (clássicos) e tokenizer próprio (LSTM/BERTimbau).", 0),
    ("Cinco classificadores, semente fixa (42), split 80/20 estratificado.", 0),
], size=16)
shrink_body(s_met2, 1.18, 1.6, 11.2, 1.35)
visoes_table(s_met2, 3.25)

# 10 RESULTADOS 1/4
decorate_title(s_r1)
set_title(s_r1, "Resultados (1/4): Domínio Sintético e Real")
set_body(s_r1, [
    ("No mundo sintético, a classificação é quase perfeita.", 0),
    ("Nos dados reais, o desempenho fica no esperado para a tarefa.", 0),
], size=16)
shrink_body(s_r1, 1.18, 1.62, 11.0, 1.0)
stat_card(s_r1, 1.7, 3.0, 4.2, 3.0, "97%", "Sintético → Sintético\n(V1, BERTimbau)")
stat_card(s_r1, 7.4, 3.0, 4.2, 3.0, "60–67%", "Dados reais\n(V2 e V4)")

# 11 RESULTADOS 2/4
decorate_title(s_r2)
set_title(s_r2, "Resultados (2/4): Cross-domain — a Queda")
set_body(s_r2, [
    ("Treino sintético, teste real: o desempenho despenca.", 0),
    ("O modelo acerta a classe dominante e erra as minoritárias.", 0),
    ("A Neutra é confundida com Positiva e Negativa.", 0),
], size=15)
shrink_body(s_r2, 1.18, 1.6, 5.5, 1.6)
stat_card(s_r2, 1.18, 3.4, 2.55, 2.7, "43%", "Filmes e séries\n(V5)")
stat_card(s_r2, 4.0, 3.4, 2.55, 2.7, "56%", "Aplicativos\n(V5)")
add_img(s_r2, os.path.join(IMG, "matriz_confusao_v3_svm_filmes.png"), 7.15, 2.0, 5.0)

# 12 RESULTADOS 3/4 (reality gap)
decorate_title(s_rgap)
set_title(s_rgap, "Resultados (3/4): Reality Gap")
set_body(s_rgap, [("A frase do LLM é limpa e direta; a review real tem gíria, erro e ironia.", 0)], size=16)
shrink_body(s_rgap, 1.18, 1.6, 11.0, 0.7)
ex = s_rgap.shapes.add_table(4, 3, Inches(0.8), Inches(2.5), Inches(11.7), Inches(3.8)).table
ex.columns[0].width = Inches(1.5)
ex.columns[1].width = Inches(5.1)
ex.columns[2].width = Inches(5.1)
exdata = [
    ["Classe", "Sintética (Claude)", "Real (UTLC-Movies)"],
    ["Positiva", "Nunca vi atuações tão intensas e verdadeiras; merece todos os prêmios.", "Estou sem xão sem ar sem vida PQP q filme"],
    ["Negativa", "A história é completamente previsível; adivinhei o final nos dez primeiros minutos.", "Horrivel com péssimas interpretações."],
    ["Neutra", "O filme tem duas horas e quinze minutos e estreou em março.", "é bem interesante! mais ainda sim não deixa de ser sessão da tarde"],
]
for ri, row in enumerate(exdata):
    for ci, val in enumerate(row):
        cell = ex.cell(ri, ci)
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        cell.margin_left = cell.margin_right = Pt(6)
        cell.margin_top = cell.margin_bottom = Pt(3)
        cell.text_frame.word_wrap = True
        cell.text_frame.clear()
        p = cell.text_frame.paragraphs[0]
        r = p.add_run()
        r.text = val
        if ri == 0:
            p.alignment = PP_ALIGN.CENTER
            r.font.bold = True
            r.font.size = Pt(14)
            r.font.color.rgb = BRANCO
            cell.fill.solid()
            cell.fill.fore_color.rgb = VINHO
        elif ci == 0:
            p.alignment = PP_ALIGN.CENTER
            r.font.bold = True
            r.font.size = Pt(13)
            r.font.color.rgb = VINHO
            cell.fill.solid()
            cell.fill.fore_color.rgb = CINZA
        else:
            r.font.size = Pt(12)
            r.font.color.rgb = PRETO
            cell.fill.solid()
            cell.fill.fore_color.rgb = BRANCO if ri % 2 else CINZA

# 13 RESULTADOS 4/4 (volume)
decorate_title(s_volume)
set_title(s_volume, "Resultados (4/4): Volume e Nichos")
set_body(s_volume, [
    ("Triplicar o volume sintético (200 → 600 frases/classe) não fecha o gap.", 0),
    ("A limitação é estrutural (distribuição), não falta de dados.", 0),
    ("Apps supera filmes e séries no desbalanceado: vocabulário mais objetivo.", 0),
], size=16)
shrink_body(s_volume, 1.18, 1.55, 11.0, 1.5)
add_img(s_volume, os.path.join(IMG, "comparativo_200_vs_600_movies.png"), 3.56, 3.2, 6.2)

# 14 CONCLUSAO 1/3
decorate_title(s_conc1)
set_title(s_conc1, "Conclusão (1/3): Objetivo Atendido")
set_body(s_conc1, [
    ("O objetivo foi atendido: avaliamos a viabilidade de treinar com dados sintéticos de LLMs.", 0),
    ("No domínio sintético, funciona muito bem (97% de F1-macro).", 0),
    ("No cross-domain realista, cai para 43% (filmes) e 56% (apps) — inviável para uso prático.", 0),
    ("A limitação é estrutural (distribuição), não falta de dados.", 0),
])

# 15 CONCLUSAO 2/3
decorate_title(s_conc2)
set_title(s_conc2, "Conclusão (2/3): Dificuldades e Limitações")
set_body(s_conc2, [
    ("Geração manual das 1.800 frases pelas interfaces web (~3–4 h por nicho).", 0),
    ("BERT estourou a RAM do Colab com o dataset real inteiro; resolvido com subamostra estratificada de 100 mil (distribuição preservada).", 0),
    ("Paradoxo de avaliação: só dá para calibrar o sintético comparando com dado real.", 0),
    ("Escopo: dois nichos; o TF-IDF não captura negação e ironia (por isso testamos LSTM e BERT).", 0),
])

# 16 CONCLUSAO 3/3
decorate_title(s_conc3)
set_title(s_conc3, "Conclusão (3/3): Trabalhos Futuros")
set_body(s_conc3, [
    ("Adaptação de domínio: combinar pouco dado real com bastante sintético para fechar o gap.", 0),
    ("Testar outros nichos com graus diferentes de subjetividade.", 0),
    ("Prompts mais elaborados e refinamento iterativo da geração.", 0),
    ("Gerar via API, com controle de temperatura e mais LLMs.", 0),
])

# 17 REFERENCIAS
decorate_title(s_refs)
set_body(s_refs, [
    ("PANG, B.; LEE, L.; VAITHYANATHAN, S. Thumbs up? Sentiment classification using machine learning techniques. EMNLP, 2002.", 0),
    ("SOUZA, F.; FILHO, J. A. Sentiment analysis on Brazilian Portuguese user reviews. IEEE LA-CCI, 2022.", 0),
    ("ARAÚJO, G.; MELO, T.; FIGUEIREDO, C. M. S. Is ChatGPT an effective solver of sentiment analysis tasks in Portuguese? PROPOR, 2024.", 0),
    ("ZHANG, W. et al. Sentiment analysis in the era of large language models: a reality check. NAACL Findings, 2024.", 0),
    ("LI, Z. et al. Synthetic data generation with large language models for text classification. EMNLP, 2023.", 0),
    ("HELLWIG, N. C.; FEHLE, J.; WOLFF, C. Exploring LLMs for the generation of synthetic training samples for aspect-based sentiment analysis. Expert Systems with Applications, 2025.", 0),
], size=14)

# 18 ENCERRAMENTO
set_title(s_fim, "Obrigado pela atenção!")
set_lines(ph(s_fim, 1), ["Davi Romauski Meurer", "davimeurer@alunos.utfpr.edu.br", "Orientador: Prof. Marlon Marcon"])

# ===================== NOTAS (fixacao) =====================
NOTAS = [
    (s_capa, "Slide 1 — Capa", """Abertura. Diz teu nome, o título do trabalho e quem é o orientador. 20–30 segundos, sem pressa. Respira e olha pra banca antes de começar.
Não esquecer: é só a abertura, não adianta resultado aqui.
Dica: "Bom dia, meu nome é Davi, vou apresentar meu TCC sobre usar LLMs para gerar dados de treino em análise de sentimento.\""""),

    (s_sum, "Slide 2 — Sumário", """Dá só o mapa do que vem: introdução, objetivos, justificativa, referencial, métodos, resultados e conclusão. Passa rápido, é pra situar a banca, não pra ler item por item.
Não esquecer: 10 segundos. Não enrola aqui."""),

    (s_intro, "Slide 3 — Introdução", """Aqui você arma o problema. A ideia: análise de sentimento é dizer se uma opinião é positiva, negativa ou neutra. Para um modelo aprender isso, precisa de muito texto rotulado, e montar isso à mão é caro e demorado — pior ainda em português, onde tem pouco dado pronto. Aí vem a sacada: e se um LLM gerar esse texto já rotulado?
Não esquecer: o problema central é o CUSTO e a escassez de dado rotulado em pt-BR. Não conta resultado ainda.
Se perguntarem "por que isso importa?": porque para o português faltam datasets, e gerar com LLM sairia muito mais barato que pagar anotação humana."""),

    (s_obj, "Slide 4 — Objetivos", """O geral é um só: descobrir se dá para treinar um classificador de sentimento usando só dados sintéticos de LLM. Os específicos são os passos: medir a assertividade no sintético, ver se os 3 LLMs concordam entre si, testar no cross-domain (treina sintético, testa real) e comparar os dois nichos.
Não esquecer: o cross-domain é o objetivo que realmente responde "funciona no mundo real?".
Se perguntarem "qual o mais importante?": o cross-domain (treino sintético → teste real), porque é o que prova utilidade prática."""),

    (s_just, "Slide 5 — Justificativa", """Por que fazer assim. Anotar dado é caro e o português tem pouco recurso, então gerar com LLM é uma alternativa atraente. Usei 3 LLMs para não depender do viés de um só. E testei 2 nichos + cross-domain justamente para ver ONDE a ideia funciona, não só se funciona.
Não esquecer: 3 LLMs = reduzir viés de fonte única; 2 nichos = ver generalização.
Se perguntarem "por que 3 LLMs e não 1?": empresas e arquiteturas diferentes (OpenAI, Google, Anthropic); reduz o viés de um gerador só e deixa o dataset mais variado."""),

    (s_refteo, "Slide 6 — Referencial Teórico", """Os conceitos que sustentam o trabalho. Análise de sentimento como classificação supervisionada (Pang, 2002, é o clássico). A vetorização TF-IDF alimentando os modelos clássicos (Naive Bayes, Regressão Logística, SVM). Os neurais: LSTM (aprende do zero) e BERTimbau (já vem pré-treinado em português). E a métrica principal, F1-macro.
Não esquecer: F1-macro porque trata todas as classes igual — não deixa a classe majoritária (Positiva) inflar o número.
Se perguntarem "por que F1-macro e não acurácia?": em dado desbalanceado, um modelo que só chuta a classe dominante tem acurácia alta mas é inútil; o F1-macro expõe isso porque pesa as 3 classes igualmente.
Se perguntarem "o que é TF-IDF?": uma forma de transformar texto em números dando peso às palavras mais distintivas de cada frase."""),

    (s_mat, "Slide 7 — Materiais", """As ferramentas. Os 3 LLMs geradores (ChatGPT, Gemini, Claude), os 5 classificadores, a stack em Python (scikit-learn, PyTorch, Transformers), os dados reais do Kaggle (UTLC) e o Colab com GPU T4 para rodar o BERT.
Não esquecer: dado sintético foi gerado pela interface web dos LLMs; dado real veio do Kaggle via kagglehub.
Se perguntarem "por que BERTimbau e não BERT normal?": o BERTimbau é treinado em português brasileiro, então entende melhor gíria e estrutura do nosso idioma."""),

    (s_met1, "Slide 8 — Métodos (1/2): Geração e Dados", """Mostra o pipeline (geração → pré-processamento → treino → avaliação) e foca na parte de dados. Foram 1.800 frases sintéticas por nicho (3 LLMs × 3 classes × 200 frases). Os prompts eram iguais nos dois nichos (isomorfos), só trocando "filme" por "app". O dado real veio do UTLC, ~100 mil reviews por nicho, e a nota de 1 a 5 virou as 3 classes (4–5 positiva, 3 neutra, 1–2 negativa).
Não esquecer: prompts isomorfos = qualquer diferença no resultado é do nicho, não do prompt.
Se perguntarem "os prompts foram iguais?": sim, idênticos entre os nichos; cada LLM interpreta diferente (ChatGPT ~50 caracteres, Gemini ~61, Claude ~95)."""),

    (s_met2, "Slide 9 — Métodos (2/2): Avaliação", """O protocolo. Pré-processamento (minúsculas, tira pontuação, remove duplicata), TF-IDF para os clássicos e tokenizer próprio para LSTM/BERT. Seed 42 para reprodutibilidade. E o coração: as 5 visões. Explica a tabela — V1 é tudo sintético, V2/V4 tudo real, V3/V5 é o cross-domain (treina sintético, testa real). V4 e V5 usam a distribuição real desbalanceada.
Não esquecer: V2 e V3 usam o MESMO teste; V4 e V5 também. Só muda a fonte do treino — assim isolo o efeito do sintético.
Se perguntarem "por que controlar o volume?": para a comparação ser justa — se o sintético tem menos dado que o real, a queda poderia ser do volume e não da fonte. Por isso a V3 iguala o volume.
Se perguntarem "por que não usar dado real no treino?": porque a pergunta é se o SINTÉTICO serve; treinar no real não responderia isso."""),

    (s_r1, "Slide 10 — Resultados (1/4): Domínio Sintético e Real", """Primeiro os dois cenários "fáceis". Quando treina e testa no próprio sintético (V1), o BERTimbau chega a 97% de F1-macro — ou seja, os dados gerados são coerentes entre si. E nos dados reais (V2 e V4), fica entre 60% e 67%, que é o normal para essa tarefa em português.
Não esquecer: 97% NÃO quer dizer que é bom no mundo real — é só sintético testando sintético. É o teto teórico.
Se perguntarem "97% não é ótimo?": é, mas só prova consistência interna do sintético; o teste de verdade é o cross-domain, no próximo slide."""),

    (s_r2, "Slide 11 — Resultados (2/4): Cross-domain — a Queda", """O resultado central. Quando treina no sintético e testa no real (V5, cenário realista desbalanceado), despenca para 43% em filmes e 56% em apps. A matriz de confusão mostra por quê: o modelo manda quase tudo para as classes extremas e erra a Neutra, que some no meio. É o viés de classe majoritária.
Não esquecer: 43% e 56% são o número que mata a viabilidade prática. A Neutra é a mais sacrificada.
Se perguntarem "por que cai tanto?": o texto sintético é regular demais, não tem a bagunça da review real (gíria, erro, ironia) — é o reality gap, que mostro no próximo slide.
Se perguntarem "por que apps vai melhor que filmes?": review de app é mais objetiva e repetitiva ("travou", "não abre"); filme tem linguagem mais subjetiva e variada."""),

    (s_rgap, "Slide 12 — Resultados (3/4): Reality Gap", """Aqui você mostra o porquê da queda de forma concreta. Compara, lado a lado, uma frase que o LLM gerou (limpinha, bem escrita) com uma review real (com gíria, sem acento, ironia). O modelo aprendeu no texto bonito e não reconhece o texto bagunçado do usuário real.
Não esquecer: esse slide explica visualmente a queda do slide anterior. Lê em voz alta a real da Positiva ("PQP q filme") — a banca entende na hora.
Se perguntarem "esse exemplo é representativo?": sim, peguei do próprio dataset real (UTLC); a tabela completa está no artigo."""),

    (s_volume, "Slide 13 — Resultados (4/4): Volume e Nichos", """Aqui respondo a pergunta óbvia: "e se gerar mais dado sintético?". Testei 200 vs 600 frases por classe (triplo). No gráfico, as barras do cross-domain (V3) ficam praticamente iguais — mais dado não fecha o gap. Isso prova que o problema é estrutural (o sintético é diferente do real), não falta de quantidade.
Não esquecer: essa é a evidência de que o problema NÃO é volume. É o argumento mais forte do trabalho.
Se perguntarem "e se gerar 10x mais?": o experimento sugere que não adianta — o gap é de distribuição, não de tamanho; precisaria mudar a forma de gerar, não a quantidade."""),

    (s_conc1, "Slide 14 — Conclusão (1/3): Objetivo Atendido", """Fecha respondendo a pergunta do trabalho. O objetivo foi atendido: dá para dizer com dado que treinar só com sintético não é viável para uso prático. Funciona lindo no sintético (97%), mas no cross-domain realista cai para 43–56%. E a limitação é estrutural, não de volume.
Não esquecer: "objetivo atendido" — a banca gosta de ouvir isso explícito. O achado é a inviabilidade prática medida com número.
Se perguntarem "então o trabalho deu negativo?": não — provar que NÃO funciona, e medir o quanto e por quê, é uma contribuição científica válida."""),

    (s_conc2, "Slide 15 — Conclusão (2/3): Dificuldades e Limitações", """Honestidade técnica — a banca valoriza muito. As dificuldades reais: gerar as 1.800 frases na mão pelas interfaces deu trabalho (3–4 h por nicho); o BERT estourou a memória do Colab com o dataset real inteiro, e resolvi pegando uma subamostra estratificada de 100 mil (mantendo a distribuição); e o paradoxo de que só dá para melhorar o sintético comparando com dado real. As limitações: só dois nichos, e o TF-IDF não pega negação nem ironia.
Não esquecer: a subamostra de 100k foi decisão de viabilidade computacional, com distribuição preservada — fala isso com segurança, não como falha.
Se perguntarem "a subamostra não enviesa?": não, é estratificada (mantém as proporções 71/20/9); é amostragem estatística válida, não corte arbitrário."""),

    (s_conc3, "Slide 16 — Conclusão (3/3): Trabalhos Futuros", """O que dá para fazer depois. O mais promissor é adaptação de domínio: misturar um pouco de dado real com bastante sintético para ver se fecha o gap. Também testar outros nichos, melhorar os prompts (refinamento iterativo) e gerar via API com controle de temperatura e mais LLMs.
Não esquecer: o futuro mais forte é o "pouco real + muito sintético" — é a saída natural do que descobri.
Se perguntarem "o que você faria diferente?": já começaria validando no real desde o início e usaria API com controle de temperatura, em vez da interface web."""),

    (s_refs, "Slide 17 — Referências", """Não leia as referências. Só deixa no ar: "estas são as principais referências do trabalho". Se a banca quiser detalhar alguma, aí você comenta.
Não esquecer: passar rápido, 5 segundos."""),

    (s_fim, "Slide 18 — Encerramento", """Agradece e abre para perguntas. "Obrigado pela atenção, estou à disposição para as perguntas."
Não esquecer: respira, não tem pressa. Se travar numa pergunta, pode dizer "boa pergunta, deixa eu pensar" — é melhor que responder errado.
Dica geral: você fala bem de improviso, então não decora frase. Sabe o número-chave de cada slide (97%, 43%, 56%, "volume não fecha o gap") e o resto sai natural."""),
]
for s, _, txt in NOTAS:
    set_notes(s, txt)

# ===================== SALVAR =====================
prs.save(OUT_ROOT)
import shutil
shutil.copyfile(OUT_ROOT, OUT_PASTA)

# Roteiro markdown
linhas = ["# Roteiro de Apresentação — TCC Davi",
          "",
          "> Material para **fixar e lembrar** (não para ler). Você fala de improviso; aqui é só para reforçar a ideia de cada slide, o número que importa e as perguntas que a banca pode fazer.",
          "",
          "**Dica geral:** saiba o número-chave de cada slide (97%, 43%, 56%, \"volume não fecha o gap\") — o resto sai natural. Tempo-alvo: ~15–18 min.",
          ""]
for _, titulo, txt in NOTAS:
    linhas.append(f"## {titulo}")
    linhas.append("")
    linhas.append(txt)
    linhas.append("")
with open(ROTEIRO, "w", encoding="utf-8") as f:
    f.write("\n".join(linhas))

print("Salvo:", OUT_ROOT)
print("Copia:", OUT_PASTA)
print("Roteiro:", ROTEIRO)
print("Total de slides:", len(prs.slides._sldIdLst))
