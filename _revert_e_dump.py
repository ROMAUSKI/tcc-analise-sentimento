"""Remove os icones de teste (ICONE_AUTO*) e imprime o texto completo dos
slides-alvo para planejar os elementos visuais. Faz backup antes de salvar."""
import glob, os, shutil, datetime
from pptx import Presentation

BASE = r"C:\Users\Davi\Documentos\tcc-analise-sentimento"
PPTX = glob.glob(os.path.join(BASE, "Apresenta*", "Apresentacao_TCC_Davi.pptx"))[0]

ts = datetime.datetime.now().strftime("%Y%m%d_%H%M%S")
bdir = os.path.join(os.path.dirname(PPTX), "_backups")
os.makedirs(bdir, exist_ok=True)
shutil.copy2(PPTX, os.path.join(bdir, f"Apresentacao_TCC_Davi_backup_{ts}.pptx"))

prs = Presentation(PPTX)
rem = 0
for slide in prs.slides:
    for sh in list(slide.shapes):
        if sh.shape_type == 13 and sh.name.startswith("ICONE_AUTO"):
            sh._element.getparent().remove(sh._element)
            rem += 1
prs.save(PPTX)
print(f"[revert] {rem} icone(s) removido(s)\n")

prs = Presentation(PPTX)
alvo = [1, 2, 3, 4, 5, 6, 13, 14, 15, 16, 17]
for idx in alvo:
    s = prs.slides[idx]
    titulo = ""
    for sh in s.shapes:
        if sh.has_text_frame and sh.name == "PlaceHolder 1":
            titulo = sh.text_frame.text
    print(f"\n##### Slide {idx+1}: {titulo!r}")
    for sh in s.shapes:
        if sh.has_text_frame and sh.name == "PlaceHolder 2":
            for p in sh.text_frame.paragraphs:
                t = "".join(r.text for r in p.runs)
                if t.strip():
                    print(f"   [n{p.level}] {t}")
