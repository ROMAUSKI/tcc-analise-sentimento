# -*- coding: utf-8 -*-
"""Monta o slide 9 (Metodos 2/3: Prompts): tabela de prompts isomorfos +
exemplo do CSV gerado. Remove os chevrons vazios do template. Backup antes.
Mexe SOMENTE no slide 9."""
import glob, os, shutil, datetime
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

BASE = r"C:\Users\Davi\Documentos\tcc-analise-sentimento"
PPTX = glob.glob(os.path.join(BASE, "Apresenta*", "Apresentacao_TCC_Davi.pptx"))[0]
VINHO = RGBColor(0x8C, 0x23, 0x32)
BRANCO = RGBColor(0xFF, 0xFF, 0xFF)
PRETO = RGBColor(0x22, 0x22, 0x22)
CINZA = RGBColor(0xEC, 0xEC, 0xEC)
SLIDE_IDX = 8  # slide 9


def label(slide, text, left, top, w, size=14):
    tb = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(w), Inches(0.34))
    p = tb.text_frame.paragraphs[0]
    r = p.add_run(); r.text = text
    r.font.size = Pt(size); r.font.bold = True; r.font.color.rgb = VINHO
    return tb


def tabela(slide, left, top, width, col_w, dados, aligns, fsize=11):
    nrows, ncols = len(dados), len(dados[0])
    gf = slide.shapes.add_table(nrows, ncols, Inches(left), Inches(top),
                                Inches(width), Inches(0.4 * nrows)).table
    for c, w in enumerate(col_w):
        gf.columns[c].width = Inches(w)
    for ri, row in enumerate(dados):
        for ci, val in enumerate(row):
            cell = gf.cell(ri, ci)
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            cell.margin_top = cell.margin_bottom = Pt(2)
            cell.margin_left = cell.margin_right = Pt(6)
            tf = cell.text_frame; tf.word_wrap = True; tf.clear()
            p = tf.paragraphs[0]; p.alignment = aligns[ci]
            r = p.add_run(); r.text = val; r.font.size = Pt(fsize)
            if ri == 0:
                r.font.bold = True; r.font.color.rgb = BRANCO
                cell.fill.solid(); cell.fill.fore_color.rgb = VINHO
            else:
                r.font.color.rgb = PRETO
                cell.fill.solid(); cell.fill.fore_color.rgb = BRANCO if ri % 2 else CINZA
    return gf


def main():
    ts = datetime.datetime.now().strftime("%Y%m%d_%H%M%S")
    bdir = os.path.join(os.path.dirname(PPTX), "_backups")
    os.makedirs(bdir, exist_ok=True)
    shutil.copy2(PPTX, os.path.join(bdir, f"Apresentacao_TCC_Davi_backup_{ts}.pptx"))

    prs = Presentation(PPTX)
    slide = prs.slides[SLIDE_IDX]

    # remove chevrons vazios do template (mantem o marcador do titulo)
    for sh in list(slide.shapes):
        if sh.name.startswith("Chevron"):
            sh._element.getparent().remove(sh._element)

    L = PP_ALIGN.LEFT; C = PP_ALIGN.CENTER
    # --- prompts ---
    label(slide, "Prompt isomorfo — mesmo para ChatGPT, Gemini e Claude (só muda a classe)", 1.18, 1.6, 11.0)
    prompts = [
        ["Classe", "Prompt enviado (resumo)"],
        ["Positiva", "Escreva 200 frases positivas avaliando filmes e séries — elogios/satisfação. Máx. 30 palavras; saída em CSV."],
        ["Negativa", "Escreva 200 frases negativas avaliando filmes e séries — críticas/insatisfação. Máx. 30 palavras; saída em CSV."],
        ["Neutra", "Escreva 200 frases neutras sobre filmes e séries — informações objetivas, sem opinião. Máx. 30 palavras; saída em CSV."],
    ]
    tabela(slide, 1.18, 1.98, 11.0, [1.7, 9.3], prompts, [C, L], fsize=11)

    # --- exemplo CSV ---
    label(slide, "Exemplo do CSV gerado", 1.18, 4.25, 11.0)
    csv = [
        ["frase", "classe", "fonte"],
        ["Este filme toca a alma de uma forma muito especial.", "Positiva", "Gemini"],
        ["Diálogos cheios de frases de efeito forçadas que ninguém diria na vida real.", "Negativa", "Claude"],
        ["A produção executiva gerencia o orçamento e cronograma das filmagens.", "Neutra", "ChatGPT"],
    ]
    tabela(slide, 1.18, 4.63, 11.0, [7.4, 1.8, 1.8], csv, [L, C, C], fsize=11)

    prs.save(PPTX)
    print(f"[salvo] {PPTX}")


if __name__ == "__main__":
    main()
