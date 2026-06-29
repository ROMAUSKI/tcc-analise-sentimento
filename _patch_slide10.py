# -*- coding: utf-8 -*-
"""Conserta o slide 10 (Estilo de Cada LLM): preenche os 2 bullets no corpo e
copia o placeholder do numero de slide (campo automatico) de um slide vizinho.
Backup antes. Mexe so no slide 10."""
import glob, os, shutil, datetime
from copy import deepcopy
from pptx import Presentation
from pptx.util import Inches, Pt

BASE = r"C:\Users\Davi\Documentos\tcc-analise-sentimento"
PPTX = glob.glob(os.path.join(BASE, "Apresenta*", "Apresentacao_TCC_Davi.pptx"))[0]

ts = datetime.datetime.now().strftime("%Y%m%d_%H%M%S")
bdir = os.path.join(os.path.dirname(PPTX), "_backups")
os.makedirs(bdir, exist_ok=True)
shutil.copy2(PPTX, os.path.join(bdir, f"Apresentacao_TCC_Davi_backup_{ts}.pptx"))

prs = Presentation(PPTX)
s = prs.slides[9]  # slide 10

# corpo = placeholder idx 0 que NAO e o titulo (name != Title*)
body = None
for sh in s.shapes:
    if sh.is_placeholder and sh.placeholder_format.idx == 0 and not sh.name.lower().startswith("title"):
        body = sh
if body is None:
    raise SystemExit("corpo nao encontrado")

body.left, body.top, body.width, body.height = Inches(1.18), Inches(1.55), Inches(11.2), Inches(1.2)
tf = body.text_frame; tf.clear(); tf.word_wrap = True
bullets = [
    "Cada LLM tem um padrão de comprimento próprio — ChatGPT mais curto (~50), Gemini no meio (~60) e Claude o mais longo (~95 em filmes).",
    "O estilo também muda por nicho: em apps, mais técnicos e diretos, os três geradores se aproximam.",
]
for i, txt in enumerate(bullets):
    p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
    p.level = 0
    r = p.add_run(); r.text = txt; r.font.size = Pt(15)
print("[ok] bullets inseridos")

# numero do slide: copiar placeholder idx 6 de um slide vizinho
ja_tem = any(sh.is_placeholder and sh.placeholder_format.idx == 6 for sh in s.shapes)
if not ja_tem:
    src = prs.slides[7]  # slide 8 tem o placeholder idx 6 (campo automatico)
    num_el = None
    for sh in src.shapes:
        if sh.is_placeholder and sh.placeholder_format.idx == 6:
            num_el = sh._element
    if num_el is not None:
        s.shapes._spTree.append(deepcopy(num_el))
        print("[ok] numero de slide copiado")
    else:
        print("[aviso] placeholder de numero nao encontrado no slide vizinho")
else:
    print("[ok] numero ja presente")

prs.save(PPTX)
print(f"[salvo] {PPTX}")
