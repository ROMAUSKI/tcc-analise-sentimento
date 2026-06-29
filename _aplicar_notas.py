# -*- coding: utf-8 -*-
"""Reescreve as notas (speaker notes) dos 20 slides em estilo FLASHCARD:
objetivo, pra bater o olho e lembrar o significado/porque. Nao e roteiro de
leitura. Regenera tambem o Roteiro_Apresentacao.md. Backup antes."""
import glob, os, shutil, datetime
from pptx import Presentation

BASE = r"C:\Users\Davi\Documentos\tcc-analise-sentimento"
PPTX = glob.glob(os.path.join(BASE, "Apresenta*", "Apresentacao_TCC_Davi.pptx"))[0]
ROTEIRO = glob.glob(os.path.join(BASE, "Apresenta*", "Roteiro_Apresentacao.md"))
ROTEIRO = ROTEIRO[0] if ROTEIRO else os.path.join(os.path.dirname(PPTX), "Roteiro_Apresentacao.md")

NOTAS = {
 0: "ABERTURA — nome, orientador (Prof. Marlon), tema.\n"
    "Pitch (1 frase): testar se LLMs servem pra gerar base de treino JÁ ROTULADA pra análise de sentimento em pt-BR — e se funciona em reviews reais.\n"
    "Respirar, olhar a banca, não correr.",
 1: "Roteiro da fala: Contexto → Método → Resultados → Conclusão.\n"
    "Só situar a banca; NÃO ler item por item.",
 2: "Análise de sentimento = classificar opinião em positiva / negativa / neutra.\n"
    "• Treinar exige base rotulada: cara e lenta; pt-BR tem pouco.\n"
    "• Ideia: o LLM gera o dado JÁ rotulado.\n"
    "Pergunta central (pipeline): gera → treina → testa no REAL → funciona?",
 3: "Geral: avaliar a viabilidade prática de treinar com dados sintéticos.\n"
    "3 perguntas: (1) assertividade do sintético; (2) coerência ENTRE os LLMs; (3) cross-domain = treino sintético, teste real.\n"
    "A 3ª (cross-domain) é a estrela — foi sugestão do Marlon.",
 4: "Por que importa: anotação humana é cara/lenta; pt-BR escasso.\n"
    "• LLM rotula sem anotação manual.\n"
    "• 3 LLMs reduzem viés de fonte única.\n"
    "• 2 nichos + cross-domain mostram ONDE funciona.",
 5: "Base teórica: análise de sentimento = classificação supervisionada (Pang, 2002).\n"
    "• Clássicos: TF-IDF + NB / LR / SVM.\n"
    "• Neurais: LSTM e BERTimbau (BERT pré-treinado em pt-BR).\n"
    "• Métrica: F1-macro (justo com classes desbalanceadas).",
 6: "O que usei: 3 LLMs geradores, 5 classificadores, 2 datasets reais (UTLC-Movies/Apps, via kagglehub).\n"
    "Stack: Python, scikit-learn, PyTorch, HuggingFace. Neurais rodaram no Colab (GPU T4).",
 7: "NÚMERO-CHAVE: 200 frases POR LLM por classe → 3 LLMs = 600/classe → 1.800/nicho.\n"
    "• 2 nichos (filmes/séries e apps).\n"
    "• Real: ~100k reviews/nicho; nota 1-5 → 3 classes (≥4 pos, ≤2 neg, =3 neu).\n"
    "Pipeline: geração → pré-proc (TF-IDF) → treino → avaliação.",
 8: "Prompt ISOMORFO: mesmo pros 3 LLMs, só muda a classe (comparação justa).\n"
    "• Pede 200 frases, máx. 30 palavras, saída em CSV.\n"
    "• CSV: frase | classe | fonte (qual LLM gerou).\n"
    "LEMBRAR: o '200' é POR LLM → 3×200 = 600 por classe.",
 9: "Estilo de geração (comprimento): ChatGPT ~50 (curto/regular); Claude ~95 em filmes (longo); Gemini ~60.\n"
    "• Em apps (texto técnico) os 3 se aproximam.\n"
    "Por que importa: frase longa/elaborada se afasta da review real → ajuda a explicar o gap e a diferença entre nichos.",
 10: "5 visões = FONTE do treino (sint/real) × TESTE (controlado/natural).\n"
     "• V1 = baseline sintético; V5 = cross-domain realista (a que importa).\n"
     "• Pares: V2 vs V3 e V4 vs V5 isolam a fonte; V1 vs V5 = reality gap.\n"
     "Métrica principal: F1-macro.",
 11: "V1 (sint→sint) = 97% com BERTimbau — sintético é ótimo 'em casa'.\n"
     "Real (V2/V4) = 60-67% — patamar normal da tarefa.\n"
     "ARMADILHA: 97% NÃO significa sucesso; é só dentro do mundo sintético.",
 12: "V5 (sint→real natural) DESPENCA: 43% filmes / 56% apps → inviável pra uso prático.\n"
     "• Matriz (V3-SVM): a classe Neutra é a mais sacrificada.\n"
     "Mensagem: o modelo treinado em sintético NÃO generaliza pro real.",
 13: "Por que cai: frase sintética é limpa e direta; review real tem gíria, erro, ironia.\n"
     "• Exemplos reais: 'PQP q filme', 'interesante', 'mais' (em vez de mas).\n"
     "Causa = diferença de DISTRIBUIÇÃO entre os textos.",
 14: "Teste de volume: triplicar o sintético (600 → 1.800 no nicho) NÃO fecha o gap em V3.\n"
     "• Logo: limitação é ESTRUTURAL, não falta de dados (argumento central!).\n"
     "• Apps > filmes no desbalanceado (vocabulário mais objetivo).\n"
     "Legenda: Reduzido = 600/nicho; Completo = 1.800/nicho. Redução = amostra aleatória por classe (seed 42).",
 15: "Objetivo ATENDIDO: avaliamos a viabilidade.\n"
     "• 97% no sintético, MAS 43-56% no cross-domain real → inviável pra produção.\n"
     "• Limitação é estrutural (distribuição), não falta de dados.",
 16: "• Geração manual: ~3-4h por nicho (interface web).\n"
     "• BERT estourou a RAM do Colab → subamostra estratificada de 100k (distribuição preservada).\n"
     "• Paradoxo: calibrar o sintético exige dado real.\n"
     "• TF-IDF não pega negação/ironia → por isso testei LSTM e BERT.",
 17: "• Adaptação de domínio (pouco real + muito sintético).\n"
     "• Outros nichos; prompts mais elaborados.\n"
     "• Gerar via API (controle de temperatura) e mais LLMs.",
 18: "Citar de passagem; não ler.\n"
     "Principais: Pang 2002 (base); Souza 2022 (sentimento pt-BR); Li 2023 / Hellwig 2025 (LLM gerando dado sintético).",
 19: "Agradecer e abrir pra perguntas.\n"
     "Mensagem final (1 frase): dado sintético de LLM serve pra protótipo/pesquisa, mas ainda NÃO substitui dado real em produção.",
}


def titulo(slide):
    for sh in slide.shapes:
        if sh.is_placeholder and sh.placeholder_format.idx == 0 and sh.has_text_frame and sh.text_frame.text.strip():
            return sh.text_frame.text.strip()
    for sh in slide.shapes:
        if sh.has_text_frame and sh.text_frame.text.strip():
            return sh.text_frame.text.strip()[:50]
    return "(sem titulo)"


def main():
    ts = datetime.datetime.now().strftime("%Y%m%d_%H%M%S")
    bdir = os.path.join(os.path.dirname(PPTX), "_backups")
    os.makedirs(bdir, exist_ok=True)
    shutil.copy2(PPTX, os.path.join(bdir, f"Apresentacao_TCC_Davi_backup_{ts}.pptx"))

    prs = Presentation(PPTX)
    titulos = []
    for idx, slide in enumerate(prs.slides):
        titulos.append(titulo(slide))
        if idx in NOTAS:
            slide.notes_slide.notes_text_frame.text = NOTAS[idx]
    prs.save(PPTX)
    print(f"[ok] notas aplicadas em {len(NOTAS)} slides")

    with open(ROTEIRO, "w", encoding="utf-8") as f:
        f.write("# Notas de fixação — Apresentação TCC (flashcards)\n\n")
        f.write("> Bater o olho e lembrar. Não é roteiro de leitura.\n\n")
        for idx in range(len(titulos)):
            f.write(f"## Slide {idx+1} — {titulos[idx]}\n\n")
            f.write(NOTAS.get(idx, "") + "\n\n")
    print(f"[ok] roteiro regenerado: {ROTEIRO}")


if __name__ == "__main__":
    main()
