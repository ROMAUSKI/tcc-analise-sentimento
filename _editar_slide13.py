"""
Edita SOMENTE o slide 13 (Resultados 4/4: Volume e Nichos) do deck:
- remove a imagem do painel 2x3 unico (Picture 7)
- insere as 5 visoes individuais + legenda num grid 3+2 (posicao inicial; o
  Davi reposiciona depois)

NAO toca em nenhum outro slide (python-pptx preserva o resto do arquivo).
Faz backup do .pptx antes de salvar.
"""
import glob, os, shutil, datetime
from pptx import Presentation
from pptx.util import Inches

BASE = r"C:\Users\Davi\Documentos\tcc-analise-sentimento"
PPTX = glob.glob(os.path.join(BASE, "Apresenta*", "Apresentacao_TCC_Davi.pptx"))[0]
IMGDIR = os.path.join(BASE, "resultados", "slide13_visoes")
SLIDE_IDX = 12  # slide 13 visivel

# --- backup ---
ts = datetime.datetime.now().strftime("%Y%m%d_%H%M%S")
bkp_dir = os.path.join(os.path.dirname(PPTX), "_backups")
os.makedirs(bkp_dir, exist_ok=True)
bkp = os.path.join(bkp_dir, f"Apresentacao_TCC_Davi_backup_{ts}.pptx")
shutil.copy2(PPTX, bkp)
print(f"[backup] {bkp}")

prs = Presentation(PPTX)
slide = prs.slides[SLIDE_IDX]

# --- remove a(s) imagem(ns) atuais do slide ---
removidas = 0
for sh in list(slide.shapes):
    if sh.shape_type == 13:  # PICTURE
        sh._element.getparent().remove(sh._element)
        removidas += 1
print(f"[remove] {removidas} imagem(ns) removida(s) do slide {SLIDE_IDX+1}")

# --- grid 3+2 (+ legenda na 6a celula) ---
W = 2.85                     # largura de cada figura (pol)
ASPECT = 3.2 / 5.0          # altura/largura do PNG => height = W*ASPECT
GAP_X = 0.45
n_cols = 3
total_w = n_cols * W + (n_cols - 1) * GAP_X
start_left = (13.333 - total_w) / 2
lefts = [start_left + c * (W + GAP_X) for c in range(n_cols)]
top1, top2 = 3.10, 5.00

posicoes = [
    ("v1.png",      lefts[0], top1),
    ("v2.png",      lefts[1], top1),
    ("v3.png",      lefts[2], top1),
    ("v4.png",      lefts[0], top2),
    ("v5.png",      lefts[1], top2),
    ("legenda.png", lefts[2], top2),
]
for nome, left, top in posicoes:
    path = os.path.join(IMGDIR, nome)
    slide.shapes.add_picture(path, Inches(left), Inches(top), width=Inches(W))
    print(f"[insere] {nome} em ({left:.2f}, {top:.2f}) w={W}")

prs.save(PPTX)
print(f"[salvo] {PPTX}")
