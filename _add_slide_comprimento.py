# -*- coding: utf-8 -*-
"""Insere um slide novo 'Métodos (3/4): Estilo de Cada LLM' (boxplot de
comprimento por LLM) entre Prompts e Avaliação, e renumera os métodos x/3 -> x/4.
Numeração do rodapé é campo automático (reajusta sozinha). Backup antes."""
import glob, os, shutil, datetime
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN

BASE = r"C:\Users\Davi\Documentos\tcc-analise-sentimento"
PPTX = glob.glob(os.path.join(BASE, "Apresenta*", "Apresentacao_TCC_Davi.pptx"))[0]
IMG = os.path.join(BASE, "resultados", "comprimento_frases_por_llm.png")
VINHO = RGBColor(0x8C, 0x23, 0x32)

RENOMEAR = {
    "Métodos (1/3): Geração e Dados": "Métodos (1/4): Geração e Dados",
    "Métodos (2/3): Prompts": "Métodos (2/4): Prompts",
    "Métodos (3/3): Avaliação": "Métodos (4/4): Avaliação",
}


def titulo_ph(slide):
    for sh in slide.shapes:
        if sh.is_placeholder and sh.placeholder_format.idx == 0:
            return sh
    return None


def main():
    ts = datetime.datetime.now().strftime("%Y%m%d_%H%M%S")
    bdir = os.path.join(os.path.dirname(PPTX), "_backups")
    os.makedirs(bdir, exist_ok=True)
    shutil.copy2(PPTX, os.path.join(bdir, f"Apresentacao_TCC_Davi_backup_{ts}.pptx"))

    prs = Presentation(PPTX)

    # 1) renumerar titulos dos metodos
    for s in prs.slides:
        t = titulo_ph(s)
        if t is None or not t.has_text_frame:
            continue
        para = t.text_frame.paragraphs[0]
        full = "".join(r.text for r in para.runs)
        if full in RENOMEAR and para.runs:
            para.runs[0].text = RENOMEAR[full]
            for r in list(para.runs)[1:]:
                r._r.getparent().remove(r._r)

    # 2) novo slide a partir do mesmo layout dos metodos
    layout = prs.slides[9].slide_layout
    novo = prs.slides.add_slide(layout)

    # titulo + marcador vinho
    tp = titulo_ph(novo)
    tp.left, tp.top, tp.width, tp.height = Inches(1.18), Inches(0.4), Inches(11.0), Inches(1.45)
    tf = tp.text_frame; tf.clear()
    r = tf.paragraphs[0].add_run(); r.text = "Métodos (3/4): Estilo de Cada LLM"
    mk = novo.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.55), Inches(0.92), Inches(0.42), Inches(0.42))
    mk.fill.solid(); mk.fill.fore_color.rgb = VINHO; mk.line.fill.background(); mk.shadow.inherit = False

    # corpo: 2 bullets de leitura
    body = None
    for ph in novo.placeholders:
        if ph.placeholder_format.idx not in (0, 6, 10):
            body = ph
    if body is not None:
        body.left, body.top, body.width, body.height = Inches(1.18), Inches(1.55), Inches(11.2), Inches(1.2)
        tfb = body.text_frame; tfb.clear(); tfb.word_wrap = True
        bullets = [
            "Cada LLM tem um padrão de comprimento próprio — ChatGPT mais curto (~50), Gemini no meio (~60) e Claude o mais longo (~95 em filmes).",
            "O estilo também muda por nicho: em apps, mais técnicos e diretos, os três geradores se aproximam.",
        ]
        for i, txt in enumerate(bullets):
            p = tfb.paragraphs[0] if i == 0 else tfb.add_paragraph()
            p.level = 0
            rr = p.add_run(); rr.text = txt; rr.font.size = Pt(15)

    # remove placeholder de data, se existir
    for ph in list(novo.placeholders):
        if ph.placeholder_format.idx == 10:
            ph._element.getparent().remove(ph._element)

    # imagem (boxplot) centralizada embaixo
    W = 10.6
    H = W * (789 / 2221)
    left = (13.333 - W) / 2
    novo.shapes.add_picture(IMG, Inches(left), Inches(2.85), width=Inches(W))

    # 3) reordenar: mover o novo (ultimo) para a posicao 9 (10o slide)
    sldIdLst = prs.slides._sldIdLst
    ids = list(sldIdLst)
    sldIdLst.remove(ids[-1])
    sldIdLst.insert(9, ids[-1])

    prs.save(PPTX)
    print(f"[salvo] {PPTX} | total slides: {len(prs.slides.__iter__.__self__._sldIdLst)}")


if __name__ == "__main__":
    main()
